import os
import json
import pdfplumber
from docx import Document
from pptx import Presentation
import tiktoken
import threading
from services.embedding_service import store_embeddings
from services.qa_service import generate_single_summary
from config import Config

def extract_text(filepath):
    ext = filepath.split('.')[-1].lower()
    if ext == 'pdf':
        with pdfplumber.open(filepath) as pdf:
            text = ''
            for page in pdf.pages:
                text += page.extract_text() or ''
    elif ext == 'docx':
        doc = Document(filepath)
        text = '\n'.join([para.text for para in doc.paragraphs])
    elif ext == 'pptx':
        prs = Presentation(filepath)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
    elif ext == 'txt':
        with open(filepath, 'r', encoding='utf-8') as f:
            text = f.read()
    else:
        raise ValueError("Unsupported file type")
    # Basic cleanup
    text = text.replace('\n\n', '\n').strip()
    return text

def chunk_text(text, chunk_size=512, overlap=100):
    """
    Optimized chunking for better retrieval quality and minimal AI usage.
    Uses larger chunks with overlap to maintain context while reducing total chunks.
    """
    encoding = tiktoken.get_encoding("cl100k_base")  # For GPT-3.5/4
    tokens = encoding.encode(text)

    if len(tokens) <= chunk_size:
        return [text]  # Return whole text if it's small enough

    chunks = []
    start = 0
    while start < len(tokens):
        end = min(start + chunk_size, len(tokens))
        chunk_tokens = tokens[start:end]
        chunk_text = encoding.decode(chunk_tokens)

        # Ensure chunks end at sentence boundaries when possible
        if end < len(tokens) and not chunk_text.endswith(('.', '!', '?', '\n')):
            # Find last sentence ending in chunk
            last_period = max(chunk_text.rfind('.'), chunk_text.rfind('!'), chunk_text.rfind('?'))
            if last_period > chunk_size * 0.7:  # Only if we're not losing too much content
                chunk_text = chunk_text[:last_period + 1]

        chunks.append(chunk_text.strip())
        start += chunk_size - overlap

        # Prevent infinite loop
        if start >= end:
            break

    return chunks

def generate_summary_background(text, filename, user_id):
    """
    Background task to generate summary for a document.
    """
    try:
        generate_single_summary(text, filename)
        print(f"Summary generated for {filename}")
    except Exception as e:
        print(f"Error generating summary for {filename}: {str(e)}")

def process_document(filepath, filename, user_id):
    """
    Process a document: extract text, chunk it, store embeddings, and generate summary in background.
    Returns a dictionary with processing results.
    """
    try:
        # Extract text from the document
        text = extract_text(filepath)

        # Chunk the text
        chunks = chunk_text(text)

        # Store embeddings
        store_embeddings(chunks, filename, user_id)

        # Save processed text to file
        processed_filepath = os.path.join('processed', f"{filename}.txt")
        os.makedirs('processed', exist_ok=True)
        with open(processed_filepath, 'w', encoding='utf-8') as f:
            f.write(text)

        # Generate summary in background thread
        summary_thread = threading.Thread(target=generate_summary_background, args=(text, filename, user_id))
        summary_thread.daemon = True
        summary_thread.start()

        return {
            'message': 'Document processed successfully',
            'chunks_count': len(chunks),
            'text_length': len(text)
        }
    except Exception as e:
        return {
            'error': f'Failed to process document: {str(e)}'
        }
